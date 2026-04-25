#!/usr/bin/env python3
"""edit_doc.py — Skill-bundled document editor for DOCX/PPTX/PDF/XLSX.

Invoked via the app's `run_shell` tool. Reads/writes files on disk;
never accepts document content over argv. Lazy-imports backend libs
so `check` works with nothing installed.

Exit codes:
  0 success
  1 generic error
  2 missing dependency
  3 bad input / path validation failed
"""
from __future__ import annotations

import argparse
import difflib
import json
import os
import shutil
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path
from typing import Iterable, Any
from xml.dom import minidom

MAX_EXTRACT_BYTES = 50 * 1024
HISTORY_DIR = ".nexa/doc-history"
EXCEL_ERRORS = ("#VALUE!", "#DIV/0!", "#REF!", "#NAME?", "#NULL!", "#NUM!", "#N/A")


# ---------------------------------------------------------------------------
# Utilities
# ---------------------------------------------------------------------------

def _die(msg: str, code: int = 1) -> None:
    print(msg, file=sys.stderr)
    sys.exit(code)


def _missing(pkg: str) -> None:
    print(f"MISSING_DEP: {pkg}", file=sys.stderr)
    print(f"Install with: python -m pip install {pkg}", file=sys.stderr)
    sys.exit(2)


def _validate_path(raw: str, must_exist: bool = True) -> Path:
    if not raw:
        _die("ERROR: --path is required", 3)
    p = Path(raw)
    if not p.is_absolute():
        _die(f"ERROR: --path must be absolute: {raw}", 3)
    try:
        resolved = p.resolve()
        cwd = Path.cwd().resolve()
        # Basic traversal guard: resolved path must live under cwd.
        resolved.relative_to(cwd)
    except ValueError:
        _die(f"ERROR: path escapes workspace: {raw}", 3)
    except OSError as e:
        _die(f"ERROR: cannot resolve path: {e}", 3)
    if must_exist and not resolved.exists():
        _die(f"ERROR: file not found: {resolved}", 3)
    return resolved


def _validate_output_path(raw: str, suffixes: set[str]) -> Path:
    p = _validate_path(raw, must_exist=False)
    if _ext(p) not in suffixes:
        _die(f"ERROR: output path must end with one of: {', '.join(sorted(suffixes))}", 3)
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def _validate_output_dir(raw: str, *, allow_existing: bool = True) -> Path:
    if not raw:
        _die("ERROR: output directory is required", 3)
    p = Path(raw)
    if not p.is_absolute():
        _die(f"ERROR: output directory must be absolute: {raw}", 3)
    try:
        resolved = p.resolve()
        resolved.relative_to(Path.cwd().resolve())
    except ValueError:
        _die(f"ERROR: output directory escapes workspace: {raw}", 3)
    except OSError as e:
        _die(f"ERROR: cannot resolve output directory: {e}", 3)
    if resolved.exists() and not allow_existing:
        _die(f"ERROR: output directory already exists: {resolved}", 3)
    resolved.mkdir(parents=True, exist_ok=True)
    return resolved


def _ext(p: Path) -> str:
    return p.suffix.lower().lstrip(".")


def _parse_pages(spec: str | None, total: int) -> list[int]:
    if not spec:
        return list(range(total))
    out: list[int] = []
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, b = part.split("-", 1)
            out.extend(range(int(a) - 1, int(b)))
        else:
            out.append(int(part) - 1)
    return [i for i in out if 0 <= i < total]


def _truncate(text: str) -> str:
    raw = text.encode("utf-8")
    if len(raw) <= MAX_EXTRACT_BYTES:
        return text
    cut = raw[:MAX_EXTRACT_BYTES].decode("utf-8", errors="ignore")
    return cut + f"\n\n[TRUNCATED: output exceeded {MAX_EXTRACT_BYTES} bytes]"


def _read_json(path: str) -> dict[str, Any]:
    spec_path = _validate_path(path)
    with spec_path.open("r", encoding="utf-8") as f:
        data = json.load(f)
    if not isinstance(data, dict):
        _die("ERROR: JSON spec root must be an object", 3)
    return data


def _read_text(path: str) -> str:
    text_path = _validate_path(path)
    return text_path.read_text(encoding="utf-8")


def _find_soffice() -> str | None:
    for name in ("soffice", "soffice.com", "libreoffice"):
        found = shutil.which(name)
        if found:
            return found
    for candidate in (
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files\LibreOffice\program\soffice.com",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ):
        if Path(candidate).exists():
            return candidate
    return None


def _find_pdftoppm() -> str | None:
    for name in ("pdftoppm", "pdftoppm.exe"):
        found = shutil.which(name)
        if found:
            return found
    return None


def _soffice_env() -> dict[str, str]:
    env = os.environ.copy()
    env.setdefault("SAL_USE_VCLPLUGIN", "svp")
    return env


def _soffice_base_cmd(soffice: str, profile_dir: Path) -> list[str]:
    return [
        soffice,
        f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
        "--headless",
        "--invisible",
        "--norestore",
        "--nolockcheck",
        "--nodefault",
    ]


def _run_soffice_convert(path: Path, to: str, outdir: Path) -> subprocess.CompletedProcess[str]:
    soffice = _find_soffice()
    if not soffice:
        _die("MISSING_DEP: LibreOffice/soffice\nInstall LibreOffice and ensure soffice is on PATH.", 2)
    with tempfile.TemporaryDirectory(prefix="nexa-lo-profile-") as profile:
        cmd = [
            *_soffice_base_cmd(soffice, Path(profile)),
            "--convert-to",
            to,
            "--outdir",
            str(outdir),
            str(path),
        ]
        return subprocess.run(
            cmd,
            text=True,
            capture_output=True,
            check=False,
            env=_soffice_env(),
        )


def _expected_converted_path(input_path: Path, outdir: Path, to: str) -> Path:
    ext = to.split(":", 1)[0].split()[0].lstrip(".")
    return outdir / f"{input_path.stem}.{ext}"


def _pretty_xml_file(path: Path) -> None:
    try:
        raw = path.read_bytes()
        parsed = minidom.parseString(raw)
        path.write_text(parsed.toprettyxml(indent="  "), encoding="utf-8")
    except Exception:
        # Not all Office XML parts are worth normalizing; keep original bytes if parsing fails.
        pass


# ---------------------------------------------------------------------------
# check
# ---------------------------------------------------------------------------

def cmd_check(_args: argparse.Namespace) -> int:
    backends = [
        ("python-docx", "docx"),
        ("python-pptx", "pptx"),
        ("pypdf", "pypdf"),
        ("openpyxl", "openpyxl"),
    ]
    missing_core = []
    print(f"python: {sys.version.split()[0]}")
    for display, mod in backends:
        try:
            imported = __import__(mod)
            ver = getattr(imported, "__version__", "unknown")
            print(f"  {display:<14} OK      ({ver})")
        except ImportError:
            print(f"  {display:<14} MISSING")
            missing_core.append(display)
        except Exception as e:  # noqa: BLE001
            # Backend present but broken (e.g. numpy ABI mismatch). Treat as missing.
            print(f"  {display:<14} BROKEN  ({type(e).__name__}: {e})")
            missing_core.append(display)
    soffice = _find_soffice()
    if soffice:
        print(f"  LibreOffice    OK      ({soffice})")
    else:
        print("  LibreOffice    MISSING (needed for convert/recalc/render QA)")
    pdftoppm = _find_pdftoppm()
    if pdftoppm:
        print(f"  Poppler        OK      ({pdftoppm})")
    else:
        print("  Poppler        MISSING (needed for render QA)")
    if missing_core:
        print(
            "\nInstall missing deps with:\n"
            f"  python -m pip install {' '.join(missing_core)}"
        )
        return 2
    return 0


# ---------------------------------------------------------------------------
# extract
# ---------------------------------------------------------------------------

def _extract_docx(path: Path) -> str:
    try:
        import docx  # type: ignore
    except ImportError:
        _missing("python-docx")
    doc = docx.Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_pptx(path: Path, pages: str | None) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        _missing("python-pptx")
    prs = Presentation(str(path))
    indices = _parse_pages(pages, len(prs.slides))
    out = []
    for i in indices:
        slide = prs.slides[i]
        out.append(f"--- Slide {i + 1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    out.append(para.text)
    return "\n".join(out)


def _extract_pdf(path: Path, pages: str | None) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError:
        _missing("pypdf")
    reader = PdfReader(str(path))
    indices = _parse_pages(pages, len(reader.pages))
    out = []
    for i in indices:
        out.append(f"--- Page {i + 1} ---")
        out.append(reader.pages[i].extract_text() or "")
    return "\n".join(out)


def _extract_xlsx(path: Path, sheets: str | None) -> str:
    try:
        import openpyxl  # type: ignore
    except ImportError:
        _missing("openpyxl")
    wb = openpyxl.load_workbook(str(path), data_only=False, read_only=True)
    wanted = {s.strip() for s in sheets.split(",")} if sheets else None
    out: list[str] = []
    for ws in wb.worksheets:
        if wanted and ws.title not in wanted:
            continue
        out.append(f"--- Sheet: {ws.title} ---")
        for row in ws.iter_rows(values_only=True):
            values = ["" if cell is None else str(cell) for cell in row]
            if any(v.strip() for v in values):
                out.append("\t".join(values).rstrip())
    return "\n".join(out)


def cmd_extract(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    ext = _ext(path)
    if ext == "docx":
        text = _extract_docx(path)
    elif ext == "pptx":
        text = _extract_pptx(path, args.pages)
    elif ext == "pdf":
        text = _extract_pdf(path, args.pages)
    elif ext == "xlsx":
        text = _extract_xlsx(path, args.sheets)
    else:
        _die(f"ERROR: extract does not support .{ext}", 3)
    sys.stdout.write(_truncate(text))
    if not text.endswith("\n"):
        sys.stdout.write("\n")
    return 0


# ---------------------------------------------------------------------------
# replace / redact (shared core)
# ---------------------------------------------------------------------------

def _iter_docx_runs(doc) -> Iterable:
    for para in doc.paragraphs:
        for run in para.runs:
            yield run
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        yield run


def _replace_docx(path: Path, find: str, replace: str, dry_run: bool) -> int:
    try:
        import docx  # type: ignore
    except ImportError:
        _missing("python-docx")
    doc = docx.Document(str(path))
    before = "\n".join(p.text for p in doc.paragraphs)
    count = 0
    for run in _iter_docx_runs(doc):
        if find in run.text:
            count += run.text.count(find)
            if not dry_run:
                run.text = run.text.replace(find, replace)
    if dry_run:
        after = before.replace(find, replace)
        diff = difflib.unified_diff(
            before.splitlines(), after.splitlines(),
            fromfile=str(path), tofile=f"{path} (preview)", lineterm="",
        )
        sys.stdout.write("\n".join(diff) + "\n")
        print(f"\n[DRY-RUN] matches: {count}")
        return 0
    doc.save(str(path))
    print(f"replaced {count} occurrence(s) in {path}")
    return 0


def _replace_pptx(path: Path, find: str, replace: str, dry_run: bool) -> int:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        _missing("python-pptx")
    prs = Presentation(str(path))
    count = 0
    before_lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    before_lines.append(run.text)
                    if find in run.text:
                        count += run.text.count(find)
                        if not dry_run:
                            run.text = run.text.replace(find, replace)
    if dry_run:
        before = "\n".join(before_lines)
        after = before.replace(find, replace)
        diff = difflib.unified_diff(
            before.splitlines(), after.splitlines(),
            fromfile=str(path), tofile=f"{path} (preview)", lineterm="",
        )
        sys.stdout.write("\n".join(diff) + "\n")
        print(f"\n[DRY-RUN] matches: {count}")
        return 0
    prs.save(str(path))
    print(f"replaced {count} occurrence(s) in {path}")
    return 0


def _replace_xlsx(path: Path, find: str, replace: str, dry_run: bool) -> int:
    try:
        import openpyxl  # type: ignore
    except ImportError:
        _missing("openpyxl")
    wb = openpyxl.load_workbook(str(path))
    before_lines: list[str] = []
    after_lines: list[str] = []
    count = 0
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and find in cell.value:
                    before = f"{ws.title}!{cell.coordinate}: {cell.value}"
                    after_value = cell.value.replace(find, replace)
                    after = f"{ws.title}!{cell.coordinate}: {after_value}"
                    before_lines.append(before)
                    after_lines.append(after)
                    count += cell.value.count(find)
                    if not dry_run:
                        cell.value = after_value
    if dry_run:
        diff = difflib.unified_diff(
            before_lines, after_lines,
            fromfile=str(path), tofile=f"{path} (preview)", lineterm="",
        )
        sys.stdout.write("\n".join(diff) + "\n")
        print(f"\n[DRY-RUN] matches: {count}")
        return 0
    wb.save(str(path))
    print(f"replaced {count} occurrence(s) in {path}")
    return 0


def cmd_replace(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    if not args.find:
        _die("ERROR: --find is required", 3)
    ext = _ext(path)
    if ext == "docx":
        return _replace_docx(path, args.find, args.replace or "", args.dry_run)
    if ext == "pptx":
        return _replace_pptx(path, args.find, args.replace or "", args.dry_run)
    if ext == "xlsx":
        return _replace_xlsx(path, args.find, args.replace or "", args.dry_run)
    _die(f"ERROR: replace supports .docx/.pptx/.xlsx only (got .{ext})", 3)
    return 1


def cmd_redact(args: argparse.Namespace) -> int:
    # redact is replace with a default mask token
    args.replace = args.replace if args.replace is not None else "[REDACTED]"
    return cmd_replace(args)


# ---------------------------------------------------------------------------
# insert_slide
# ---------------------------------------------------------------------------

def cmd_insert_slide(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    if _ext(path) != "pptx":
        _die("ERROR: insert_slide requires a .pptx file", 3)
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches  # type: ignore
    except ImportError:
        _missing("python-pptx")
    prs = Presentation(str(path))
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    # Populate title/body if placeholders exist.
    if slide.shapes.title is not None:
        slide.shapes.title.text = args.title or ""
    if args.body:
        body_placeholder = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                body_placeholder = ph
                break
        if body_placeholder is None:
            left = top = Inches(1)
            width = Inches(8)
            height = Inches(5)
            body_placeholder = slide.shapes.add_textbox(left, top, width, height)
        body_placeholder.text_frame.text = args.body

    # Reorder: move new slide to position after --after.
    after = max(0, int(args.after))
    xml_slides = prs.slides._sldIdLst  # noqa: SLF001
    slides = list(xml_slides)
    new_el = slides[-1]
    xml_slides.remove(new_el)
    xml_slides.insert(after, new_el)

    prs.save(str(path))
    print(f"inserted slide after position {after} in {path}")
    return 0


# ---------------------------------------------------------------------------
# version
# ---------------------------------------------------------------------------

def cmd_version(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    root = Path.cwd() / HISTORY_DIR / path.name
    root.mkdir(parents=True, exist_ok=True)
    existing = sorted(
        (p for p in root.iterdir() if p.is_dir() and p.name.startswith("v")),
        key=lambda p: int(p.name[1:]) if p.name[1:].isdigit() else 0,
    )
    next_n = 1
    if existing:
        last = existing[-1].name
        try:
            next_n = int(last[1:]) + 1
        except ValueError:
            next_n = len(existing) + 1
    dest_dir = root / f"v{next_n}"
    dest_dir.mkdir()
    dest = dest_dir / path.name
    shutil.copy2(path, dest)
    print(f"v{next_n} -> {dest}")
    return 0


# ---------------------------------------------------------------------------
# create_docx / create_xlsx / create_pptx
# ---------------------------------------------------------------------------

def _numbered_markdown_text(stripped: str) -> str | None:
    prefix, sep, rest = stripped.partition(".")
    if sep and prefix.isdigit() and rest.startswith(" "):
        return rest.strip()
    return None


def _docx_add_markdown(doc, markdown: str) -> None:
    lines = markdown.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        if not stripped:
            doc.add_paragraph()
            i += 1
            continue
        if stripped.startswith("#"):
            level = len(stripped) - len(stripped.lstrip("#"))
            if 1 <= level <= 6 and len(stripped) > level and stripped[level].isspace():
                doc.add_heading(stripped[level:].strip(), level=min(level, 4))
                i += 1
                continue
        if stripped.startswith("|") and stripped.endswith("|"):
            table_lines: list[str] = []
            while i < len(lines) and lines[i].strip().startswith("|") and lines[i].strip().endswith("|"):
                table_lines.append(lines[i].strip())
                i += 1
            rows = [[cell.strip() for cell in row.strip("|").split("|")] for row in table_lines]
            rows = [
                row for row in rows
                if not all(cell and set(cell) <= {"-", ":", " "} for cell in row)
            ]
            if rows:
                table = doc.add_table(rows=len(rows), cols=max(len(r) for r in rows))
                table.style = "Table Grid"
                for ri, row in enumerate(rows):
                    for ci, value in enumerate(row):
                        table.rows[ri].cells[ci].text = value
                        if ri == 0:
                            for paragraph in table.rows[ri].cells[ci].paragraphs:
                                for run in paragraph.runs:
                                    run.bold = True
                continue
        numbered = _numbered_markdown_text(stripped)
        if stripped.startswith(("- ", "* ", "• ")):
            doc.add_paragraph(stripped[2:].strip(), style="List Bullet")
        elif numbered is not None:
            doc.add_paragraph(numbered, style="List Number")
        elif stripped.startswith("> "):
            p = doc.add_paragraph(stripped[2:].strip())
            p.style = "Intense Quote"
        else:
            doc.add_paragraph(stripped)
        i += 1


def cmd_create_docx(args: argparse.Namespace) -> int:
    try:
        import docx  # type: ignore
    except ImportError:
        _missing("python-docx")
    path = _validate_output_path(args.path, {"docx"})
    doc = docx.Document(str(_validate_path(args.template))) if args.template else docx.Document()
    if args.font:
        doc.styles["Normal"].font.name = args.font
    if args.title:
        doc.add_heading(args.title, 0)
    if args.subtitle:
        p = doc.add_paragraph(args.subtitle)
        p.style = "Subtitle" if "Subtitle" in [s.name for s in doc.styles] else p.style
    if args.input_md:
        _docx_add_markdown(doc, _read_text(args.input_md))
    elif args.body:
        _docx_add_markdown(doc, args.body)
    if args.footer:
        for section in doc.sections:
            section.footer.paragraphs[0].text = args.footer
    if args.author:
        doc.core_properties.author = args.author
    doc.save(str(path))
    print(f"created DOCX: {path}")
    return 0


def cmd_create_xlsx(args: argparse.Namespace) -> int:
    try:
        import openpyxl  # type: ignore
        from openpyxl.chart import BarChart, LineChart, Reference  # type: ignore
        from openpyxl.styles import Font, PatternFill, Alignment  # type: ignore
        from openpyxl.worksheet.table import Table, TableStyleInfo  # type: ignore
    except ImportError:
        _missing("openpyxl")
    path = _validate_output_path(args.path, {"xlsx"})
    spec = _read_json(args.spec)
    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    sheets = spec.get("sheets") or []
    if not isinstance(sheets, list) or not sheets:
        _die("ERROR: create_xlsx spec requires non-empty 'sheets' array", 3)
    for sheet_spec in sheets:
        name = str(sheet_spec.get("name") or f"Sheet{len(wb.worksheets) + 1}")[:31]
        ws = wb.create_sheet(name)
        rows = sheet_spec.get("rows") or []
        for row in rows:
            ws.append(row if isinstance(row, list) else [row])
        if rows:
            for cell in ws[1]:
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill("solid", fgColor="2B579A")
                cell.alignment = Alignment(horizontal="center")
            ws.freeze_panes = sheet_spec.get("freeze_panes") or "A2"
            ws.auto_filter.ref = ws.dimensions
            if sheet_spec.get("table", True):
                table_name = "".join(ch for ch in name if ch.isalnum()) or f"Table{len(wb.worksheets)}"
                table = Table(displayName=f"{table_name[:24]}Table", ref=ws.dimensions)
                style = TableStyleInfo(
                    name="TableStyleMedium2",
                    showFirstColumn=False,
                    showLastColumn=False,
                    showRowStripes=True,
                    showColumnStripes=False,
                )
                table.tableStyleInfo = style
                ws.add_table(table)
        for item in sheet_spec.get("formulas") or []:
            ws[str(item["cell"])] = str(item["formula"])
        for idx, width in enumerate(sheet_spec.get("column_widths") or [], start=1):
            ws.column_dimensions[openpyxl.utils.get_column_letter(idx)].width = float(width)
        for chart_spec in sheet_spec.get("charts") or []:
            chart_type = str(chart_spec.get("type") or "bar").lower()
            chart = LineChart() if chart_type == "line" else BarChart()
            chart.title = chart_spec.get("title") or ""
            min_col = int(chart_spec.get("min_col", 2))
            max_col = int(chart_spec.get("max_col", min_col))
            min_row = int(chart_spec.get("min_row", 1))
            max_row = int(chart_spec.get("max_row", ws.max_row))
            data = Reference(ws, min_col=min_col, max_col=max_col, min_row=min_row, max_row=max_row)
            chart.add_data(data, titles_from_data=True)
            if chart_spec.get("categories_col"):
                cats = Reference(ws, min_col=int(chart_spec["categories_col"]), min_row=min_row + 1, max_row=max_row)
                chart.set_categories(cats)
            ws.add_chart(chart, chart_spec.get("anchor") or "E2")
    wb.save(str(path))
    print(f"created XLSX: {path}")
    return 0


def cmd_create_pptx(args: argparse.Namespace) -> int:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
    except ImportError:
        _missing("python-pptx")
    path = _validate_output_path(args.path, {"pptx"})
    spec = _read_json(args.spec)
    prs = Presentation(str(_validate_path(args.template))) if args.template else Presentation()
    # Remove default empty slide only if present and unused by templates.
    if not args.template and len(prs.slides) == 0:
        pass
    slides = spec.get("slides") or []
    if not isinstance(slides, list) or not slides:
        _die("ERROR: create_pptx spec requires non-empty 'slides' array", 3)
    for slide_spec in slides:
        layout_name = str(slide_spec.get("layout") or "body").lower()
        layout_idx = 0 if layout_name == "title" else 1 if len(prs.slide_layouts) > 1 else 0
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        if slide.shapes.title is not None:
            slide.shapes.title.text = str(slide_spec.get("title") or "")
        subtitle = slide_spec.get("subtitle")
        body = slide_spec.get("body")
        bullets = slide_spec.get("bullets") or []
        if subtitle and layout_name == "title" and len(slide.placeholders) > 1:
            slide.placeholders[1].text = str(subtitle)
        elif body or bullets:
            box = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    box = ph
                    break
            if box is None:
                box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.7), Inches(4.8))
            tf = box.text_frame
            tf.clear()
            if body:
                tf.paragraphs[0].text = str(body)
                tf.paragraphs[0].font.size = Pt(18)
            for item in bullets:
                p = tf.add_paragraph()
                p.text = str(item)
                p.level = 0
                p.font.size = Pt(18)
        table = slide_spec.get("table")
        if isinstance(table, list) and table:
            rows = len(table)
            cols = max(len(row) for row in table if isinstance(row, list))
            shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(2.0), Inches(8.7), Inches(3.6))
            for ri, row in enumerate(table):
                for ci, value in enumerate(row):
                    shape.table.cell(ri, ci).text = str(value)
        notes = slide_spec.get("notes")
        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = str(notes)
            except Exception:
                pass
    prs.save(str(path))
    print(f"created PPTX: {path}")
    return 0


# ---------------------------------------------------------------------------
# ooxml / render / recalc / validate / convert
# ---------------------------------------------------------------------------

def cmd_unpack(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    if _ext(path) not in {"docx", "pptx", "xlsx"}:
        _die("ERROR: unpack supports .docx/.pptx/.xlsx only", 3)
    outdir = _validate_output_dir(args.outdir, allow_existing=True)
    if any(outdir.iterdir()):
        if not args.overwrite:
            _die(f"ERROR: output directory is not empty: {outdir}. Pass --overwrite to replace it.", 3)
        # Safety guard: _validate_output_dir already proved the resolved path is under cwd.
        shutil.rmtree(outdir)
        outdir.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(path) as zf:
        zf.extractall(outdir)
    xml_count = 0
    for member in list(outdir.rglob("*.xml")) + list(outdir.rglob("*.rels")):
        _pretty_xml_file(member)
        xml_count += 1
    print(f"unpacked {path.name} -> {outdir} ({xml_count} XML parts prettified)")
    return 0


def cmd_pack(args: argparse.Namespace) -> int:
    input_dir = _validate_path(args.input_dir)
    if not input_dir.is_dir():
        _die(f"ERROR: input directory is not a directory: {input_dir}", 3)
    if not (input_dir / "[Content_Types].xml").exists():
        _die("ERROR: input directory does not look like an unpacked Office document", 3)
    path = _validate_output_path(args.path, {"docx", "pptx", "xlsx"})
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
        for item in sorted(input_dir.rglob("*")):
            if item.is_file():
                zf.write(item, item.relative_to(input_dir).as_posix())
    with zipfile.ZipFile(path) as zf:
        bad = zf.testzip()
        if bad:
            _die(f"ERROR: corrupt Office ZIP member after pack: {bad}", 1)
    print(f"packed {input_dir} -> {path}")
    return 0


def _convert_to_pdf(path: Path, outdir: Path) -> Path:
    completed = _run_soffice_convert(path, "pdf", outdir)
    if completed.returncode != 0:
        stderr = completed.stderr.strip() or completed.stdout.strip() or "LibreOffice conversion failed"
        _die(stderr, completed.returncode or 1)
    pdf = _expected_converted_path(path, outdir, "pdf")
    if not pdf.exists():
        matches = sorted(outdir.glob("*.pdf"))
        if matches:
            return matches[0]
        _die(f"ERROR: LibreOffice did not produce a PDF in {outdir}", 1)
    return pdf


def cmd_render(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    if _ext(path) not in {"docx", "pptx", "xlsx", "pdf"}:
        _die("ERROR: render supports .docx/.pptx/.xlsx/.pdf only", 3)
    pdftoppm = _find_pdftoppm()
    if not pdftoppm:
        _die("MISSING_DEP: Poppler/pdftoppm\nInstall Poppler and ensure pdftoppm is on PATH.", 2)
    outdir = _validate_output_dir(args.outdir, allow_existing=True)
    image_format = args.format.lower()
    if image_format not in {"png", "jpeg"}:
        _die("ERROR: --format must be png or jpeg", 3)
    with tempfile.TemporaryDirectory(prefix="nexa-render-") as tmp:
        pdf = path if _ext(path) == "pdf" else _convert_to_pdf(path, Path(tmp))
        prefix = outdir / "page"
        cmd = [
            pdftoppm,
            f"-{image_format}",
            "-r",
            str(args.dpi),
            str(pdf),
            str(prefix),
        ]
        completed = subprocess.run(cmd, text=True, capture_output=True, check=False)
    if completed.stdout:
        print(completed.stdout.strip())
    if completed.stderr:
        print(completed.stderr.strip(), file=sys.stderr)
    if completed.returncode != 0:
        return completed.returncode
    images = sorted(outdir.glob(f"page*.{'jpg' if image_format == 'jpeg' else 'png'}"))
    print(f"rendered {len(images)} page image(s) to {outdir}")
    for image in images[:20]:
        print(image)
    if len(images) > 20:
        print(f"... {len(images) - 20} more")
    return 0


def _scan_xlsx_formula_errors(path: Path) -> tuple[int, dict[str, list[str]]]:
    try:
        import openpyxl  # type: ignore
    except ImportError:
        _missing("openpyxl")
    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    found: dict[str, list[str]] = {err: [] for err in EXCEL_ERRORS}
    try:
        for sheet_name in wb.sheetnames:
            for row in wb[sheet_name].iter_rows():
                for cell in row:
                    if isinstance(cell.value, str):
                        for err in EXCEL_ERRORS:
                            if err in cell.value:
                                found[err].append(f"{sheet_name}!{cell.coordinate}")
                                break
    finally:
        wb.close()
    total = sum(len(locations) for locations in found.values())
    return total, {err: locs for err, locs in found.items() if locs}


def _count_xlsx_formulas(path: Path) -> int:
    try:
        import openpyxl  # type: ignore
    except ImportError:
        _missing("openpyxl")
    wb = openpyxl.load_workbook(str(path), data_only=False, read_only=True)
    count = 0
    try:
        for sheet_name in wb.sheetnames:
            for row in wb[sheet_name].iter_rows():
                for cell in row:
                    if isinstance(cell.value, str) and cell.value.startswith("="):
                        count += 1
    finally:
        wb.close()
    return count


def cmd_recalc_xlsx(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    if _ext(path) != "xlsx":
        _die("ERROR: recalc_xlsx requires a .xlsx file", 3)
    soffice = _find_soffice()
    if not soffice:
        _die("MISSING_DEP: LibreOffice/soffice\nInstall LibreOffice and ensure soffice is on PATH.", 2)
    with tempfile.TemporaryDirectory(prefix="nexa-xlsx-recalc-") as tmp:
        tmp_path = Path(tmp)
        source = tmp_path / path.name
        outdir = tmp_path / "out"
        outdir.mkdir()
        shutil.copy2(path, source)
        completed = _run_soffice_convert(source, "xlsx", outdir)
        if completed.returncode != 0:
            stderr = completed.stderr.strip() or completed.stdout.strip() or "LibreOffice recalculation failed"
            _die(stderr, completed.returncode or 1)
        output = _expected_converted_path(source, outdir, "xlsx")
        if not output.exists():
            matches = sorted(outdir.glob("*.xlsx"))
            if not matches:
                _die("ERROR: LibreOffice did not produce a recalculated XLSX", 1)
            output = matches[0]
        shutil.copy2(output, path)
    total_errors, errors = _scan_xlsx_formula_errors(path)
    result = {
        "status": "success" if total_errors == 0 else "errors_found",
        "total_errors": total_errors,
        "total_formulas": _count_xlsx_formulas(path),
        "error_summary": {
            err: {"count": len(locations), "locations": locations[:20]}
            for err, locations in errors.items()
        },
    }
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0 if total_errors == 0 else 1

def cmd_validate(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    ext = _ext(path)
    if ext in {"docx", "pptx", "xlsx"}:
        with zipfile.ZipFile(path) as zf:
            bad = zf.testzip()
            if bad:
                _die(f"ERROR: corrupt Office ZIP member: {bad}", 1)
    if ext == "docx":
        try:
            import docx  # type: ignore
        except ImportError:
            _missing("python-docx")
        doc = docx.Document(str(path))
        print(f"VALID DOCX paragraphs={len(doc.paragraphs)} tables={len(doc.tables)}")
    elif ext == "pptx":
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            _missing("python-pptx")
        prs = Presentation(str(path))
        print(f"VALID PPTX slides={len(prs.slides)}")
    elif ext == "xlsx":
        try:
            import openpyxl  # type: ignore
        except ImportError:
            _missing("openpyxl")
        wb = openpyxl.load_workbook(str(path), data_only=False, read_only=True)
        sheet_count = len(wb.worksheets)
        sheet_names = ",".join(wb.sheetnames)
        wb.close()
        total_errors, errors = _scan_xlsx_formula_errors(path)
        formula_count = _count_xlsx_formulas(path)
        print(f"VALID XLSX sheets={sheet_count} names={sheet_names} formulas={formula_count} formula_errors={total_errors}")
        if errors:
            print(json.dumps(errors, ensure_ascii=False, indent=2))
            return 1
    elif ext == "pdf":
        try:
            from pypdf import PdfReader  # type: ignore
        except ImportError:
            _missing("pypdf")
        reader = PdfReader(str(path))
        print(f"VALID PDF pages={len(reader.pages)}")
    else:
        _die(f"ERROR: validate does not support .{ext}", 3)
    return 0


def cmd_convert(args: argparse.Namespace) -> int:
    path = _validate_path(args.path)
    outdir = Path(args.outdir).resolve() if args.outdir else path.parent
    try:
        outdir.relative_to(Path.cwd().resolve())
    except ValueError:
        _die(f"ERROR: --outdir escapes workspace: {outdir}", 3)
    outdir.mkdir(parents=True, exist_ok=True)
    completed = _run_soffice_convert(path, args.to, outdir)
    if completed.stdout:
        print(completed.stdout.strip())
    if completed.stderr:
        print(completed.stderr.strip(), file=sys.stderr)
    if completed.returncode != 0:
        return completed.returncode
    print(f"converted {path.name} -> .{args.to} in {outdir}")
    return 0


# ---------------------------------------------------------------------------
# argparse
# ---------------------------------------------------------------------------

def build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        prog="edit_doc.py",
        description="Edit existing DOCX/PPTX/PDF/XLSX documents from run_shell.",
    )
    p.add_argument("--path", help="Absolute path to the target file")
    sub = p.add_subparsers(dest="cmd", required=True)

    sub.add_parser("check", help="Report available backends").set_defaults(func=cmd_check)

    p_rep = sub.add_parser("replace", help="Replace text (docx/pptx/xlsx)")
    p_rep.add_argument("--find", required=True)
    p_rep.add_argument("--replace", default="")
    p_rep.add_argument("--dry-run", action="store_true")
    p_rep.set_defaults(func=cmd_replace)

    p_red = sub.add_parser("redact", help="Redact text (docx/pptx/xlsx)")
    p_red.add_argument("--find", required=True)
    p_red.add_argument("--replace", default=None)
    p_red.add_argument("--dry-run", action="store_true")
    p_red.set_defaults(func=cmd_redact)

    p_ext = sub.add_parser("extract", help="Extract plain text (docx/pdf/pptx/xlsx)")
    p_ext.add_argument("--pages", default=None, help="e.g. 1-3 or 1,3,5")
    p_ext.add_argument("--sheets", default=None, help="Comma-separated XLSX sheet names")
    p_ext.set_defaults(func=cmd_extract)

    p_ins = sub.add_parser("insert_slide", help="Insert a slide into a pptx")
    p_ins.add_argument("--after", type=int, default=0)
    p_ins.add_argument("--title", default="")
    p_ins.add_argument("--body", default="")
    p_ins.set_defaults(func=cmd_insert_slide)

    p_ver = sub.add_parser("version", help="Snapshot file to .nexa/doc-history")
    p_ver.set_defaults(func=cmd_version)

    p_cd = sub.add_parser("create_docx", help="Create a DOCX using python-docx")
    p_cd.add_argument("--title", default="")
    p_cd.add_argument("--subtitle", default="")
    p_cd.add_argument("--body", default="")
    p_cd.add_argument("--input-md", default=None, help="Absolute path to markdown source")
    p_cd.add_argument("--template", default=None, help="Optional absolute .docx template path")
    p_cd.add_argument("--font", default="Calibri")
    p_cd.add_argument("--footer", default="")
    p_cd.add_argument("--author", default="Nexa")
    p_cd.set_defaults(func=cmd_create_docx)

    p_cx = sub.add_parser("create_xlsx", help="Create an XLSX workbook from a JSON spec")
    p_cx.add_argument("--spec", required=True, help="Absolute path to workbook JSON spec")
    p_cx.set_defaults(func=cmd_create_xlsx)

    p_cp = sub.add_parser("create_pptx", help="Create a PPTX presentation from a JSON spec")
    p_cp.add_argument("--spec", required=True, help="Absolute path to deck JSON spec")
    p_cp.add_argument("--template", default=None, help="Optional absolute .pptx template path")
    p_cp.set_defaults(func=cmd_create_pptx)

    p_unpack = sub.add_parser("unpack", help="Unpack DOCX/PPTX/XLSX into editable OOXML")
    p_unpack.add_argument("--outdir", required=True, help="Absolute output directory")
    p_unpack.add_argument("--overwrite", action="store_true")
    p_unpack.set_defaults(func=cmd_unpack)

    p_pack = sub.add_parser("pack", help="Pack an unpacked OOXML directory back into DOCX/PPTX/XLSX")
    p_pack.add_argument("--input-dir", required=True, help="Absolute unpacked OOXML directory")
    p_pack.set_defaults(func=cmd_pack)

    p_render = sub.add_parser("render", help="Render DOCX/PPTX/XLSX/PDF pages or slides to images for visual QA")
    p_render.add_argument("--outdir", required=True, help="Absolute output directory for page images")
    p_render.add_argument("--dpi", type=int, default=150)
    p_render.add_argument("--format", default="png", choices=["png", "jpeg"])
    p_render.set_defaults(func=cmd_render)

    p_recalc = sub.add_parser("recalc_xlsx", help="Recalculate XLSX formulas with LibreOffice and scan for formula errors")
    p_recalc.set_defaults(func=cmd_recalc_xlsx)

    p_val = sub.add_parser("validate", help="Validate that a document opens with its backend")
    p_val.set_defaults(func=cmd_validate)

    p_conv = sub.add_parser("convert", help="Convert via LibreOffice headless")
    p_conv.add_argument("--to", required=True, help="Output extension/filter, e.g. pdf, docx, xlsx")
    p_conv.add_argument("--outdir", default=None, help="Optional absolute output directory")
    p_conv.set_defaults(func=cmd_convert)

    return p


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    try:
        return args.func(args)
    except SystemExit:
        raise
    except Exception as e:  # noqa: BLE001
        print(f"ERROR: {e}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    sys.exit(main())
